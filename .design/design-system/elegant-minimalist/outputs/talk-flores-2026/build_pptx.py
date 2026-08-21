#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Build an editable PPTX of the Flores 2026 talk in the elegant-minimalist style."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# palette
PAGE=RGBColor(0xF8,0xF6,0xF1); SURF=RGBColor(0xFF,0xFF,0xFF); MUTEDBG=RGBColor(0xF1,0xEE,0xE7)
TEXT=RGBColor(0x1F,0x1F,0x1B); MUT=RGBColor(0x5F,0x62,0x5C); BORDER=RGBColor(0xD8,0xD2,0xC5)
ACC=RGBColor(0x8A,0x2D,0x1F); SUP=RGBColor(0x2F,0x5D,0x62)
SERIF="Source Serif 4"; UI="Inter"
IMG_MAP="assets/w01_seismisitas_indonesia.png"
IMG_WAVE="assets/waveform_mmri_flores2026.png"
IMG_HIST="assets/wichmann_historical_map.png"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
SW,SH=prs.slide_width,prs.slide_height

def slide(bg=PAGE):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background()
    r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element)
    return s

def tb(s,x,y,w,h,anchor=MSO_ANCHOR.TOP):
    box=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=box.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    return box,tf

def run(p,txt,size,color,font=UI,bold=False,italic=False):
    r=p.add_run(); r.text=txt; f=r.font
    f.size=Pt(size); f.color.rgb=color; f.name=font; f.bold=bold; f.italic=italic
    return r

def para(tf,new=True):
    p=tf.paragraphs[0] if (not new and len(tf.paragraphs)==1 and not tf.paragraphs[0].runs) else tf.add_paragraph()
    return p

def eyebrow(s,txt,x=0.62,y=0.5,color=SUP):
    _,tf=tb(s,x,y,11,0.4); r=run(tf.paragraphs[0],txt.upper(),13,color,UI,bold=True)
    tf.paragraphs[0].runs[0].font._rPr.set('spc','900')
    return tf

def footer(s,left,idx):
    _,tf=tb(s,0.62,6.95,9,0.4); run(tf.paragraphs[0],left,10.5,MUT,UI)
    _,tf2=tb(s,11.8,6.95,1.0,0.4); p=tf2.paragraphs[0]; p.alignment=PP_ALIGN.RIGHT
    run(p,f"{idx} / 11",10.5,MUT,UI)

def rect(s,x,y,w,h,fill=None,line=None,lw=1.0,dash=None):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.shadow.inherit=False
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(lw)
    return sh

def callout(s,x,y,w,h,title,body,accent=SUP):
    rect(s,x,y,w,h,fill=MUTEDBG)
    rect(s,x,y,0.06,h,fill=accent)
    _,tf=tb(s,x+0.28,y+0.22,w-0.5,h-0.4,anchor=MSO_ANCHOR.TOP)
    p=tf.paragraphs[0]; run(p,title.upper(),12,accent,UI,bold=True); p.space_after=Pt(6)
    p2=tf.add_paragraph();
    for seg,bold in body: run(p2,seg,16,TEXT,UI,bold=bold)
    p2.line_spacing=1.25

def bullets(s,x,y,w,items,size=17,gap=10,mark=ACC):
    _,tf=tb(s,x,y,w,5)
    first=True
    for segs in items:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.space_after=Pt(gap); p.line_spacing=1.2
        run(p,"■  ",size-3,mark,UI,bold=True)
        for seg,bold in segs: run(p,seg,size,TEXT,UI,bold=bold)

def bigstat(s,x,y,num,label,numsize=48):
    _,tf=tb(s,x,y,4.5,1.5)
    p=tf.paragraphs[0]; run(p,num,numsize,ACC,SERIF,bold=True)
    p2=tf.add_paragraph(); run(p2,label,15,MUT,UI); p2.space_before=Pt(4)

def pic_cropped(s,path,x,y,w,h,cl=0,cr=0,ct=0,cb=0):
    pic=s.shapes.add_picture(path,Inches(x),Inches(y),Inches(w),Inches(h))
    pic.crop_left=cl; pic.crop_right=cr; pic.crop_top=ct; pic.crop_bottom=cb
    return pic

# ---------- 1 title ----------
s=slide()
eyebrow(s,"Diskusi Publik · PSBA UGM · 21 Agustus 2026")
_,tf=tb(s,0.6,1.35,12,2.6)
p=tf.paragraphs[0]; run(p,"Indonesia adalah",60,TEXT,SERIF,bold=True); p.line_spacing=1.0
p2=tf.add_paragraph(); run(p2,"rumah gempa.",60,ACC,SERIF,bold=True); p2.line_spacing=1.0
_,tf=tb(s,0.62,3.9,10.5,1.4)
p=tf.paragraphs[0]; p.line_spacing=1.3
run(p,"217.801 gempa dalam 26 tahun. Pertanyaannya bukan ",22,MUT,SERIF)
run(p,"apakah",22,MUT,SERIF,italic=True); run(p," akan terjadi, tapi ",22,MUT,SERIF)
run(p,"di mana, kapan,",22,MUT,SERIF,italic=True); run(p," dan ",22,MUT,SERIF)
run(p,"seberapa siap",22,MUT,SERIF,italic=True); run(p," kita.",22,MUT,SERIF)
footer(s,"Dr.rer.nat. Wiwit Suryanto — Geoscience Research Group, Departemen Fisika, FMIPA UGM",1)

# ---------- 2 full map ----------
s=slide()
eyebrow(s,"Satu peta, satu pertanyaan")
_,tf=tb(s,0.6,0.95,12.1,0.7); run(tf.paragraphs[0],"Setiap titik adalah satu gempa yang pernah mengguncang negeri ini",26,TEXT,SERIF,bold=True)
rect(s,0.6,1.75,12.13,4.75,fill=SURF,line=BORDER,lw=1)
s.shapes.add_picture(IMG_MAP,Inches(0.75),Inches(1.9),height=Inches(4.2))
_,tf=tb(s,0.75,6.15,11.5,0.4); run(tf.paragraphs[0],"Seismisitas Indonesia — katalog BMKG 1998–2024 (217.801 gempa). Warna = kedalaman hiposenter (0–700 km).",12.5,MUT,UI)
footer(s,"Sumber gambar: katalog BMKG · Kuliah Seismologi UGM",2)

# ---------- 3 zoom ----------
s=slide()
# cropped/zoomed map fills slide
pic_cropped(s,IMG_MAP,0,0,13.333,7.5,cl=0.36,cr=0.16,ct=0.34,cb=0.14)
# ring near center
ov=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(6.05),Inches(3.15),Inches(1.25),Inches(1.25))
ov.fill.background(); ov.line.color.rgb=ACC; ov.line.width=Pt(3.5); ov.shadow.inherit=False
# label
lbl=rect(s,7.5,3.55,2.7,0.5,fill=RGBColor(0xF8,0xF6,0xF1),line=BORDER,lw=1)
_,tf=tb(s,7.6,3.62,2.5,0.4); run(tf.paragraphs[0],"Flores · 15 Agu 2026",15,ACC,UI,bold=True)
# heading on translucent-ish plate
plate=rect(s,0.5,0.5,6.4,1.9,fill=RGBColor(0xF8,0xF6,0xF1))
plate.fill.fore_color.rgb=RGBColor(0xF8,0xF6,0xF1)
_,tf=tb(s,0.7,0.62,6.0,0.4); r=run(tf.paragraphs[0],"PERBESAR: NUSA TENGGARA TIMUR",13,ACC,UI,bold=True)
_,tf=tb(s,0.7,1.05,6.1,1.3); p=tf.paragraphs[0]; p.line_spacing=1.08
run(p,"Titik ini bukan kebetulan — ia bagian dari denyut kegempaan Indonesia",25,TEXT,SERIF,bold=True)
footer(s,"Busur Sunda–Banda · sesar naik busur-belakang Flores",3)

# ---------- 4 facts ----------
s=slide()
eyebrow(s,"Gempa Flores — 15 Agustus 2026")
_,tf=tb(s,0.6,0.95,8,0.7); run(tf.paragraphs[0],"Apa yang terjadi, dalam angka",26,TEXT,SERIF,bold=True)
facts=[("Waktu asal","05:58:21 WITA"),("Magnitudo","Mw 7,7 (USGS) · 7,61 (GFZ)"),
 ("Kedalaman","10 km — dangkal"),("Episenter","±68 km BBL Ende, lepas pantai utara Flores"),
 ("Mekanisme","Sesar naik (thrust)"),("Sumber","Flores Back-arc Thrust"),
 ("Guncangan maks","MMI VIII · PGA 0,635 g (Kuwus)")]
y=1.95
for k,v in facts:
    _,tf=tb(s,0.62,y,3.0,0.5); run(tf.paragraphs[0],k,15,MUT,UI,bold=False)
    _,tf=tb(s,3.7,y,3.4,0.6); run(tf.paragraphs[0],v,16,TEXT,UI,bold=True)
    rect(s,0.62,y+0.52,6.5,0.014,fill=BORDER)
    y+=0.62
bigstat(s,8.4,1.9,"±73","korban jiwa · 1.182 luka · ±59.600 mengungsi",54)
bigstat(s,8.4,3.5,">2.768","gempa susulan hingga 19 Agu (terbesar M6,2)",54)
bigstat(s,8.4,5.1,"2′16″","peringatan dini tsunami terbit setelah gempa",54)
footer(s,"Sumber: BMKG · USGS · GFZ (data masih berkembang)",4)

# ---------- 5 timeline ----------
s=slide()
eyebrow(s,"Sesar yang sama, berulang")
_,tf=tb(s,0.6,0.95,12,0.7); run(tf.paragraphs[0],"Sesar naik busur-belakang Flores: sudah lama kita kenal",26,TEXT,SERIF,bold=True)
tl=[("1992",[("Flores, Mw 7,8",True),(" — tsunami mematikan, ±2.500 korban. Pelajaran pahit pertama.",False)]),
    ("2018",[("Lombok, seri Mw 6,3–6,9",True),(" — segmen barat sesar yang sama. Data & analisis b-value kami.",False)]),
    ("2026",[("Flores, Mw 7,7",True),(" — pecah makin ke timur. Peristiwa yang kita bahas malam ini.",False)])]
y=2.1
for yr,segs in tl:
    _,tf=tb(s,0.62,y,1.3,0.6); run(tf.paragraphs[0],yr,26,ACC,UI,bold=True)
    _,tf=tb(s,2.1,y+0.03,5.0,1.0); p=tf.paragraphs[0]; p.line_spacing=1.2
    for seg,b in segs: run(p,seg,17,TEXT,UI,bold=b)
    rect(s,0.62,y+0.95,6.6,0.014,fill=BORDER)
    y+=1.25
callout(s,7.7,2.2,5.0,2.6,"Pesan kunci",
  [("Ini bukan ancaman baru. Sesar naik di utara busur Nusa Tenggara berulang kali melepaskan gempa besar. Yang bisa kita ubah bukan gempanya — tapi ",False),("kesiapan",True),(" kita.",False)])
footer(s,"Analog seismotektonik — Flores Back-arc Thrust",5)

# ---------- 6 historical (Wichmann) ----------
s=slide()
eyebrow(s,"Sebelum era seismograf")
_,tf=tb(s,0.6,0.95,12,0.7); run(tf.paragraphs[0],"Sejarah sudah mencatatnya — jauh sebelum kita bisa mengukurnya",25,TEXT,SERIF,bold=True)
hist=[("1648",[("Solor–Larantuka",True),(" — rangkaian gempa merusak paling awal tercatat di NTT; tembok Benteng Henricus runtuh.",False)]),
      ("1837",[("Laut utara Flores",True),(" — gempa laut (seaquake) kuat dirasakan dari perahu.",False)]),
      ("1855",[("Nanga Rama, Manggarai",True),(" — gelombang besar menerjang teluk (kandidat tsunami).",False)]),
      ("1868",[("Ambugaga, Ende",True),(" — guncangan sangat keras ±2 menit di Flores tengah.",False)])]
y=1.95
for yr,segs in hist:
    _,tf=tb(s,0.62,y,1.2,0.6); run(tf.paragraphs[0],yr,23,ACC,UI,bold=True)
    _,tf=tb(s,1.95,y+0.02,5.2,0.95); p=tf.paragraphs[0]; p.line_spacing=1.12
    for seg,b in segs: run(p,seg,14.5,TEXT,UI,bold=b)
    rect(s,0.62,y+0.9,6.75,0.014,fill=BORDER)
    y+=1.0
_,tf=tb(s,0.62,y+0.02,6.8,0.5); run(tf.paragraphs[0],"Catatan makroseismik historis — belum dapat diberi magnitudo atau dipastikan dari sesar yang sama.",11,MUT,UI,italic=True)
callout(s,7.7,2.2,5.0,3.2,"Benang merah",
  [("Katalog Arthur Wichmann (hingga 1877) mencatat gempa, gempa laut, dan tsunami di Flores–NTT jauh sebelum ada seismograf. ",False),("Sejarah sudah memperingatkan; 2026 menegaskannya.",True),(" Ke depan, kesiapsiagaan bukan pilihan.",False)])
footer(s,"Sumber: A. Wichmann, katalog gempa Kepulauan Hindia (hingga 1877)",6)

# ---------- 7 historical map (Wichmann) ----------
s=slide()
eyebrow(s,"Peta gempa historis — katalog Wichmann")
_,tf=tb(s,0.6,0.95,12.1,0.7); run(tf.paragraphs[0],"Jejak gempa & tsunami NTT tersebar luas — jauh sebelum 1877",25,TEXT,SERIF,bold=True)
rect(s,0.6,1.75,12.13,4.75,fill=SURF,line=BORDER,lw=1)
_pic=s.shapes.add_picture(IMG_HIST,Inches(0.75),Inches(1.9),height=Inches(4.2))
_pic.left=Inches(0.6)+(Inches(12.13)-_pic.width)//2
_,tf=tb(s,0.75,6.15,11.5,0.4); run(tf.paragraphs[0],"Lokasi perkiraan kejadian utama hingga 1877 — gempa darat/pesisir, seaquake/tsunami, peristiwa regional. Sumber: katalog A. Wichmann.",12.5,MUT,UI)
footer(s,"Sumber: A. Wichmann, katalog gempa Kepulauan Hindia (hingga 1877)",7)

# ---------- 8 72 hours ----------
s=slide()
eyebrow(s,"72 jam pertama")
_,tf=tb(s,0.6,0.95,12,0.7); run(tf.paragraphs[0],"Apa yang dikerjakan seismologi saat tanggap darurat",26,TEXT,SERIF,bold=True)
bullets(s,0.62,1.95,12,[
 [("Menit pertama — membaca “telegram” bumi. ",True),("Gelombang P tiba lebih dulu, membawa lokasi, kedalaman, magnitudo, dan mekanisme dalam hitungan menit.",False)],
 [("ShakeMap. ",True),("Dari satu sumber, kita petakan sebaran guncangan — menebak di mana kerusakan terparah sebelum laporan masuk, untuk mengarahkan SAR.",False)],
 [("Peringatan tsunami. ",True),("Keputusan besar diambil di menit-menit awal (lihat: 2 menit 16 detik).",False)],
 [("Prakiraan gempa susulan. ",True),("Melindungi penyintas dan relawan yang masuk ke bangunan retak.",False)],
],size=17,gap=12)
callout(s,0.62,5.0,12.1,1.5,"Kecepatan vs ketelitian",
 [("Mw 7,7 (USGS) vs 7,61 (GFZ); 10 km vs 25 km. Magnitudo ",False),("direvisi",True),(" seiring data masuk — ini normal, bukan kesalahan. Publik perlu memahaminya.",False)])
footer(s,"Peran geofisika dalam rantai tanggap darurat",8)

# ---------- 7 waveform ----------
s=slide()
eyebrow(s,"Rekaman langsung — GE.MMRI, Maumere",color=ACC)
_,tf=tb(s,0.6,0.95,12.2,1.0); p=tf.paragraphs[0]; p.line_spacing=1.05
run(p,"Guncangan yang meruntuhkan Pelabuhan Maumere, terekam ±104 km dari sumber",25,TEXT,SERIF,bold=True)
rect(s,0.6,2.0,7.35,4.5,fill=SURF,line=BORDER,lw=1)
s.shapes.add_picture(IMG_WAVE,Inches(0.72),Inches(2.12),width=Inches(7.1))
_,tf=tb(s,0.72,6.05,7.1,0.4); run(tf.paragraphs[0],"GE.MMRI (8,64°S 122,24°E), broadband + akselerometer, Δ≈104 km. Data: GEOFON · ObsPy.",12,MUT,UI)
# right column
_,tf=tb(s,8.35,2.15,2.3,1.2); p=tf.paragraphs[0]; run(p,"0,14 g",40,ACC,SERIF,bold=True)
p2=tf.add_paragraph(); run(p2,"PGA (akselerometer)",14,MUT,UI)
_,tf=tb(s,10.8,2.15,2.3,1.2); p=tf.paragraphs[0]; run(p,"1,9 cm/s",40,ACC,SERIF,bold=True)
p2=tf.add_paragraph(); run(p2,"PGV puncak",14,MUT,UI)
bullets(s,8.35,3.85,4.6,[
 [("Jeda waktu asal → gelombang P (~15 dtk): ",True),("gelombang butuh waktu menempuh 104 km.",False)],
 [("Rekaman ini “menyaksikan” runtuhnya terminal Pelabuhan Maumere",True),(" — jembatan ke dampak struktural.",False)],
],size=16,gap=12,mark=SUP)
footer(s,"Data: GEOFON (GFZ) · workflow ObsPy (adaptasi load.py)",9)

# ---------- 8 aftershocks + tsunami ----------
s=slide()
eyebrow(s,"Dua pertanyaan yang paling ditanya warga")
_,tf=tb(s,0.6,1.4,5.9,0.8); run(tf.paragraphs[0],"“Akan ada susulan yang lebih besar?”",22,TEXT,SERIF,bold=True)
bullets(s,0.62,2.4,5.9,[
 [(">2.768 susulan",True),(" mengikuti pola peluruhan (hukum Omori): rapat di awal, mereda seiring waktu.",False)],
 [("Peluang susulan besar mengecil tiap hari — tapi ",False),("tak pernah nol",True),(". Jangan masuk bangunan retak.",False)],
],size=17,gap=12)
callout(s,7.0,1.9,5.7,2.4,"Pesan yang menyelamatkan",
 [("Untuk tsunami dekat-pantai, ",False),("guncangan kuat itu sendiri adalah peringatan.",True),(" Jangan tunggu sirene — ",False),("segera evakuasi mandiri",True),(" ke tempat tinggi.",False)],accent=ACC)
_,tf=tb(s,7.0,4.5,5.7,0.8); run(tf.paragraphs[0],"1992: tsunami Flores tiba dalam hitungan menit. Waktu adalah nyawa.",15,MUT,UI,italic=True)
footer(s,"Komunikasi risiko saat darurat",10)

# ---------- 9 closing ----------
s=slide()
eyebrow(s,"Yang bisa kita perbaiki")
_,tf=tb(s,0.6,1.15,12,1.0); p=tf.paragraphs[0]
run(p,"Indonesia timur butuh ",36,TEXT,SERIF,bold=True); run(p,"lebih banyak mata",36,ACC,SERIF,bold=True)
bullets(s,0.62,2.4,7.0,[
 [("Kerapatan stasiun di NTT jauh di bawah Jawa → lokasi lebih tak pasti, respons lebih lambat.",False)],
 [("Seismometer biaya rendah & jaringan berbasis komunitas–sekolah",True),(" — merapatkan pemantauan sekaligus membangun kesadaran.",False)],
 [("Penyebaran alat sementara pasca-gempa untuk memantau susulan & memetakan sesar aktif.",False)],
],size=17,gap=14)
callout(s,8.0,2.4,4.7,2.6,"Penutup",
 [("Kita tak bisa mencegah gempa. Tapi setiap sensor tambahan dan setiap warga yang tahu harus ",False),("segera lari ke tempat tinggi",True),(" memperbesar peluang kita selamat saat gempa berikutnya datang.",False)])
_,tf=tb(s,0.62,5.6,12,0.6); p=tf.paragraphs[0]
run(p,"Dr.rer.nat. Wiwit Suryanto",17,TEXT,UI,bold=True)
run(p,"  ·  Geoscience Research Group, Departemen Fisika, FMIPA UGM  ·  Terima kasih.",17,MUT,UI)
footer(s,"Diskusi Publik PSBA UGM — Respons Tanggap Darurat Gempa NTT",11)

out="Diskusi_Publik_Gempa_Flores_2026_Wiwit_Suryanto.pptx"
prs.save(out)
print("Saved",out,"slides:",len(prs.slides._sldIdLst))
